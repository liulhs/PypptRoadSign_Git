from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.util import Inches
from PypptRoadSign.background import Background

BLACK = RGBColor(0, 0, 0)  # 黑色
WHITE = RGBColor(255, 255, 255)  # 白色
NAVYBLUE = RGBColor(0, 63, 143)  # 预定的蓝色
ROUND_DEGREE = 0.05  # 圆角程度调整 0-1


class Sign:
    def __init__(self, background_style='default', width=16, height=9):
        prs = Presentation()
        prs.slide_width = Inches(width)  # Slide宽度 16 Inches, 高度 9 Inches 16:9
        prs.slide_height = Inches(height)
        self.width = width
        self.height = height
        blank_slide_layout = prs.slide_layouts[6]  # index 6 是空白slide预设
        slide = prs.slides.add_slide(blank_slide_layout)  # 添加slide
        self.__prs__ = prs
        self.canvas = slide.shapes
        self.background = Background(background_style=background_style, sign=self)
        self.background.draw()

    def output(self, dir='.'):
        self.__prs__.save(dir)
